#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation # pip install python-pptx
from datetime import datetime


def buildSlide(img_FileName):
    
    prs = Presentation('./slides/template.pptx')
        
    # ========================================================
    # Add the title slide
    # ========================================================
    
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Hello QRadar & Python World!"
    subtitle.text = "Bjoern found magic on " + datetime.now().strftime('%Y-%m-%d %H:%M:%S')
   
    # ========================================================
    # Add 2nd slide with the image
    # ========================================================
    blank_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(blank_slide_layout)
    title = slide.shapes.title    
    title.text = "This is the Image Slide"
    
    picture = slide.shapes.add_picture(img_FileName,1,1)
    
    # ========================================================
    # Center the image
    picture.left = int((prs.slide_width - picture.width) / 2)
    picture.top = int((prs.slide_height - picture.height) / 2)
    
    # ========================================================
    # Add page number
    
    # ========================================================
    # Add the Thank You Slide
    # ========================================================
    thank_you_slide_layout = prs.slide_layouts[14]
    slide = prs.slides.add_slide(thank_you_slide_layout)
    
    prs.save('./slides/rpt_EventStats.pptx')


